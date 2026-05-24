"""
生成答辩演示PPT - 智能交通流量监测与预测系统
总时长: 3分钟, 8页幻灯片
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色方案 ──────────────────────────────────
C_PRIMARY   = RGBColor(0x1A, 0x52, 0x76)  # 深蓝
C_ACCENT    = RGBColor(0x29, 0x80, 0xB9)  # 中蓝
C_LIGHT_BG  = RGBColor(0xEC, 0xF0, 0xF1)  # 浅灰
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK      = RGBColor(0x2C, 0x3E, 0x50)  # 正文深灰
C_GRAY      = RGBColor(0x7F, 0x8C, 0x8D)  # 次要灰
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)  # 指标绿
C_RED       = RGBColor(0xE7, 0x4C, 0x3C)  # 强调红
C_ORANGE    = RGBColor(0xE6, 0x7E, 0x22)  # 橙色

BASE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(BASE, "frontend", "src", "assets")

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ── 工具函数 ──────────────────────────────────
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=C_DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=C_DARK, bold_first=False, line_spacing=1.5, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 0.8)
        if bold_first and i == 0:
            p.font.bold = True
    return tf

def add_rect(slide, left, top, width, height, fill_color=C_PRIMARY, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=C_WHITE, border_color=C_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape

def add_page_number(slide, num):
    add_textbox(slide, 12.3, 7.0, 0.8, 0.4, str(num), font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)

def add_top_bar(slide):
    """每页顶部蓝色装饰条"""
    add_rect(slide, 0, 0, 13.333, 0.06, C_PRIMARY)

def add_section_title(slide, title, subtitle=None):
    """统一章节标题样式"""
    add_top_bar(slide)
    add_textbox(slide, 0.8, 0.4, 11.5, 0.7, title, font_size=32, color=C_PRIMARY, bold=True)
    # 标题下方装饰线
    add_rect(slide, 0.8, 1.05, 1.5, 0.04, C_ACCENT)
    if subtitle:
        add_textbox(slide, 0.8, 1.2, 11.5, 0.5, subtitle, font_size=14, color=C_GRAY)

def add_table(slide, left, top, col_widths, headers, rows, font_size=12):
    """添加样式化表格"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(left), Inches(top), Inches(total_w), Inches(0.5 * num_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = C_DARK
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
    return table_shape

print("Functions ready, building slides...")

# ═══════════════════════════════════════════════
# 第1页：封面
# ═══════════════════════════════════════════════
slide1 = add_blank_slide()
# 深蓝背景
add_rect(slide1, 0, 0, 13.333, 7.5, C_PRIMARY)
# 装饰条
add_rect(slide1, 0, 3.15, 13.333, 0.03, C_ACCENT)
add_rect(slide1, 0, 3.25, 13.333, 0.005, C_ACCENT)
# Logo (如果有)
logo_path = os.path.join(ASSETS, "traffic.png")
if os.path.exists(logo_path):
    slide1.shapes.add_picture(logo_path, Inches(6.0), Inches(0.6), Inches(1.333), Inches(1.333))
# 论文题目
add_textbox(slide1, 1.5, 2.2, 10.333, 0.9,
    "智能交通流量监测与预测系统", font_size=40, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
# 英文副标题
add_textbox(slide1, 1.5, 2.85, 10.333, 0.4,
    "Intelligent Traffic Flow Monitoring and Prediction System", font_size=16, color=C_ACCENT, align=PP_ALIGN.CENTER)
# 个人信息区
add_textbox(slide1, 1.5, 3.7, 10.333, 0.4,
    "学号：2022XXXXXXXX  |  姓名：XXX  |  指导老师：XXX", font_size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
# 底部
add_textbox(slide1, 1.5, 5.5, 10.333, 0.5,
    "XX大学 XX学院", font_size=20, color=C_LIGHT_BG, align=PP_ALIGN.CENTER)
add_textbox(slide1, 1.5, 6.2, 10.333, 0.4,
    "2026年5月", font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)
add_page_number(slide1, 1)

# ═══════════════════════════════════════════════
# 第2页：目录
# ═══════════════════════════════════════════════
slide2 = add_blank_slide()
add_top_bar(slide2)
add_textbox(slide2, 0.8, 0.4, 11.5, 0.7, "目  录", font_size=32, color=C_PRIMARY, bold=True)
add_rect(slide2, 0.8, 1.05, 1.5, 0.04, C_ACCENT)

toc_items = [
    ("1", "研究背景与意义", "为什么做这个系统？"),
    ("2", "系统架构与功能模块", "系统是什么样的？"),
    ("3", "核心算法：LST-GCN", "怎么做到的？"),
    ("4", "测试结果与性能评估", "效果怎么样？"),
    ("5", "总结与展望", "优缺点与后续方向"),
]
for i, (num, title, desc) in enumerate(toc_items):
    y = 1.6 + i * 1.05
    add_rounded_rect(slide2, 1.2, y, 10.8, 0.85, C_WHITE if i % 2 == 0 else C_LIGHT_BG)
    # 序号圆圈
    circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(y + 0.15), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER
    # 标题
    add_textbox(slide2, 2.4, y + 0.12, 6.0, 0.45, title, font_size=22, color=C_PRIMARY, bold=True)
    add_textbox(slide2, 2.4, y + 0.48, 6.0, 0.35, desc, font_size=13, color=C_GRAY)

add_page_number(slide2, 2)

# ═══════════════════════════════════════════════
# 第3页：研究背景与意义
# ═══════════════════════════════════════════════
slide3 = add_blank_slide()
add_section_title(slide3, "1. 研究背景与意义")

# 左侧：背景
add_textbox(slide3, 0.8, 1.5, 5.5, 0.4, "▎研究背景", font_size=20, color=C_PRIMARY, bold=True)
bg_lines = [
    "● 城市交通拥堵问题日益严峻，传统固定检测设备部署成本高、覆盖盲区大",
    "● 互联网路况数据（如高德API）为低成本、广覆盖的交通感知提供了新路径",
    "● 现有交通预测方法多为单点时间序列模型，忽略了路网的空间拓扑关联",
    "● 时空图神经网络（ST-GNN）为「时间+空间」联合建模提供了技术可能",
]
add_multiline(slide3, 0.8, 2.1, 5.5, 3.2, bg_lines, font_size=14, color=C_DARK, line_spacing=2.0)

# 右侧：意义
add_textbox(slide3, 7.0, 1.5, 5.5, 0.4, "▎研究意义", font_size=20, color=C_PRIMARY, bold=True)
sig_lines = [
    "● 理论意义：提出LSTM嵌入GCN的融合架构，验证多时域直接监督策略",
    "● 工程意义：构建完整的「数据采集→存储→推理→可视化」全链路系统",
    "● 应用价值：面向成都11个真实路口，提供可演示、可部署的交通监测方案",
]
add_multiline(slide3, 7.0, 2.1, 5.5, 2.5, sig_lines, font_size=14, color=C_DARK, line_spacing=2.2)

# 底部时间线
add_textbox(slide3, 0.8, 5.2, 11.5, 0.4, "▎项目时间线", font_size=20, color=C_PRIMARY, bold=True)
add_rect(slide3, 0.8, 5.65, 11.7, 1.5, C_LIGHT_BG)
tl_text = (
    "2025.10  确定题目    2025.11-12  文献调研与开题    2026.2-3  系统框架搭建    "
    "2026.3-5  实习期间开发    2026.5  收尾部署与论文撰写"
)
add_textbox(slide3, 1.0, 5.85, 11.3, 0.4, tl_text, font_size=12, color=C_DARK)
add_textbox(slide3, 1.0, 6.25, 11.3, 0.5,
    "· 9天真实数据采集周期（2026.4.26-5.5）  · 4月中旬更换设备并整体迁移项目  · 模型训练由云端GPU迁移至本地",
    font_size=11, color=C_GRAY)
add_page_number(slide3, 3)

# ═══════════════════════════════════════════════
# 第4页：系统架构与功能模块
# ═══════════════════════════════════════════════
slide4 = add_blank_slide()
add_section_title(slide4, "2. 系统架构与功能模块", "React + Node.js/Express + Flask/PyTorch + MySQL + Redis")

# 架构图 - 用矩形表示
layers = [
    ("浏览器", C_GRAY, 1.5, 1.8),
    ("React 前端\n(Vite + 高德地图SDK)", C_ACCENT, 4.0, 1.8),
    ("Node.js/Express 后端\n(REST API + 中间件)", C_PRIMARY, 6.5, 1.8),
    ("MySQL 8\n(业务数据)", RGBColor(0x16, 0xA0, 0x85), 6.5, 4.2),
    ("Redis 6\n(缓存加速)", RGBColor(0xDC, 0x35, 0x45), 6.5, 5.2),
    ("Flask AI 推理服务\n(LST-GCN 前向传播)", RGBColor(0x8E, 0x44, 0xAD), 9.5, 4.2),
    ("Python 采集器\n(高德路况API)", C_ORANGE, 9.5, 5.2),
]
for label, color, left, top in layers:
    shape = add_rounded_rect(slide4, left, top, 2.2, 1.1, color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

# 箭头标注
arrows = [
    ("→", 3.8, 2.1),
    ("→", 6.2, 2.1),
    ("→", 6.2, 4.5),
    ("→", 6.2, 5.5),
    ("↕", 8.8, 4.5),
]
for txt, x, y in arrows:
    add_textbox(slide4, x, y, 0.5, 0.4, txt, font_size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

# 右侧：功能模块列表
add_textbox(slide4, 0.8, 6.0, 11.5, 0.35,
    "▎六大核心功能模块",
    font_size=18, color=C_PRIMARY, bold=True)

modules_text = (
    "Dashboard监控台  ·  实时路网地图  ·  突发事件管理  ·  智能路径推荐  ·  AI预测服务  ·  用户权限系统"
)
add_textbox(slide4, 0.8, 6.45, 11.5, 0.35, modules_text, font_size=14, color=C_DARK, align=PP_ALIGN.CENTER)
add_page_number(slide4, 4)

# ═══════════════════════════════════════════════
# 第5页：核心功能模块详情
# ═══════════════════════════════════════════════
slide5 = add_blank_slide()
add_section_title(slide5, "2. 核心功能模块详情")

modules = [
    ("监测总览面板", "monitor.png",
     ["11路口实时路况一站展示", "单节点全天速度曲线 + 预测对比",
      "早/午/晚高峰一键聚焦", "实测值与预测值即时拟合评估"]),
    ("实时路网地图", "map.png",
     ["高德底图 + 11节点拥堵覆盖物", "拥堵颜色编码（畅行/缓行/拥堵/严重拥堵）",
      "点击节点弹出实时速度信息窗", "与总览面板同源数据，状态同步"]),
    ("突发事件监控", "console.png",
     ["事件上报、分配、处理、忽略全流程", "管理员/执行者双层角色权限",
      "事件关联路口一键跳转地图", "前端按钮 + 后端中间件双重鉴权"]),
    ("智能路径推荐", "recommend.png",
     ["多时域（15/30/45/60分钟）速度预测对比", "速度绝对值 + 变化趋势综合评分",
      "多节点并排对比，辅助路网级调度", "自然语言通行建议（建议通行/谨慎/绕行）"]),
]

for i, (title, img_name, points) in enumerate(modules):
    x = 0.5 + (i % 4) * 3.15
    y = 1.5
    # 卡片背景
    add_rounded_rect(slide5, x, y, 2.95, 5.5, C_WHITE, C_LIGHT_BG)
    # 模块图标
    img_path = os.path.join(ASSETS, img_name)
    if os.path.exists(img_path):
        try:
            slide5.shapes.add_picture(img_path, Inches(x + 1.1), Inches(y + 0.15), Inches(0.7), Inches(0.7))
        except:
            pass
    # 模块名
    add_textbox(slide5, x + 0.15, y + 1.05, 2.65, 0.35, title, font_size=17, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide5, x + 0.5, y + 1.42, 1.95, 0.015, C_ACCENT)
    # 功能点
    for j, pt in enumerate(points):
        add_textbox(slide5, x + 0.2, y + 1.7 + j * 0.75, 2.55, 0.65,
            f"• {pt}", font_size=12, color=C_DARK)
add_page_number(slide5, 5)

# ═══════════════════════════════════════════════
# 第6页：核心算法 LST-GCN
# ═══════════════════════════════════════════════
slide6 = add_blank_slide()
add_section_title(slide6, "3. 核心算法：LST-GCN 时空融合模型")

# 左侧：创新点
add_textbox(slide6, 0.8, 1.5, 5.5, 0.4, "▎核心创新：LSTM 嵌入 GCN 消息传递", font_size=18, color=C_PRIMARY, bold=True)
innovation = [
    "传统做法：GCN → LSTM 串联堆叠，时间建模与空间建模分离",
    "",
    "本系统做法：将LSTM单元直接嵌入GCN的消息传递过程——",
    "每个时间步的节点特征更新由LSTM单元完成，当前时间步的邻居",
    "聚合特征与上一时间步的LSTM隐状态共同输入LSTM，输出新隐状",
    "态作为节点更新后的特征表示。",
    "",
    "效果：时间建模与空间建模在同一计算步骤内同步完成，避免了",
    "串行堆叠带来的信息损耗。",
]
add_multiline(slide6, 0.8, 2.1, 5.5, 4.5, innovation, font_size=13, color=C_DARK, line_spacing=1.4)

# 右侧：技术规格
add_textbox(slide6, 7.0, 1.5, 5.5, 0.4, "▎模型关键参数", font_size=18, color=C_PRIMARY, bold=True)

params_headers = ["参数", "取值"]
params_rows = [
    ["输入窗口", "12 × 5min（60分钟历史）"],
    ["GCN隐藏层维度", "64"],
    ["LSTM隐状态维度", "64"],
    ["GCN层数", "2（覆盖二阶邻居）"],
    ["多时域输出", "4（15/30/45/60分钟）"],
    ["邻接矩阵", "对称归一化（含自环）"],
    ["Dropout", "0.1"],
]
add_table(slide6, 7.0, 2.1, [1.8, 3.5], params_headers, params_rows, font_size=12)

# 多时域直接监督策略
add_textbox(slide6, 7.0, 5.6, 5.5, 0.4, "▎多时域直接监督策略", font_size=18, color=C_PRIMARY, bold=True)
strategy = [
    "四个预测头的输入共享同一套时空特征提取层参数，",
    "仅在最终全连接映射层各自使用独立的线性变换。",
    "避免了单步滚动预测中误差随步长指数级堆积的问题，",
    "兼顾参数效率与各时域预测的独立性。",
]
add_multiline(slide6, 7.0, 6.05, 5.5, 1.3, strategy, font_size=12, color=C_DARK, line_spacing=1.5)
add_page_number(slide6, 6)

# ═══════════════════════════════════════════════
# 第7页：测试结果与性能评估
# ═══════════════════════════════════════════════
slide7 = add_blank_slide()
add_section_title(slide7, "4. 测试结果与性能评估")

# 模型指标表
add_textbox(slide7, 0.8, 1.5, 5.8, 0.4, "▎LST-GCN 多时域预测精度（测试集后20%）", font_size=17, color=C_PRIMARY, bold=True)
eval_headers = ["预测时域", "MAE (km/h)", "RMSE (km/h)", "MAPE"]
eval_rows = [
    ["15分钟", "3.18", "4.62", "8.73%"],
    ["30分钟", "3.37", "4.82", "9.14%"],
    ["45分钟", "3.82", "5.20", "10.42%"],
    ["60分钟", "3.64", "5.11", "10.02%"],
]
add_table(slide7, 0.8, 2.05, [1.3, 1.5, 1.5, 1.3], eval_headers, eval_rows, font_size=13)

# 关键结论
eval_conclusions = [
    "✓ 四时域MAE均 < 4 km/h，MAPE均 < 11%",
    "✓ 多时域最大最小MAPE差距仅1.7%，无误差指数堆积",
    "✓ 60分钟误差略低于45分钟（目标时刻速度波动更平缓）",
    "✓ 训练/验证损失曲线走势吻合，无过拟合",
]
add_multiline(slide7, 0.8, 4.35, 5.8, 2.5, eval_conclusions, font_size=13, color=C_DARK, line_spacing=2.0)

# 性能指标表
add_textbox(slide7, 7.0, 1.5, 5.5, 0.4, "▎核心接口响应时间（云端CPU环境）", font_size=17, color=C_PRIMARY, bold=True)
perf_headers = ["接口", "响应时间"]
perf_rows = [
    ["/api/traffic/latest", "18ms"],
    ["/api/predict/latest", "21ms"],
    ["/api/health", "28ms"],
    ["/api/incidents", "41ms"],
    ["/api/dashboard/chart", "100ms"],
    ["/api/report/predict-export", "170ms"],
    ["/api/route/outlook", "208ms"],
    ["/api/predict/trigger", "326ms"],
]
add_table(slide7, 7.0, 2.05, [2.8, 1.6], perf_headers, perf_rows, font_size=12)

# 性能结论
perf_conclusions = [
    "✓ 普通查询接口 < 50ms（Redis缓存 < 20ms）",
    "✓ 推理类接口（trigger/outlook）< 400ms",
    "✓ 前端首次加载：Dashboard 428ms, Map 2030ms (含SDK)",
    "✓ 16项功能测试用例全部通过",
    "✓ 三层架构云端持续运行一周无异常",
]
add_multiline(slide7, 7.0, 5.75, 5.5, 1.6, perf_conclusions, font_size=12, color=C_DARK, line_spacing=1.8)
add_page_number(slide7, 7)

# ═══════════════════════════════════════════════
# 第8页：总结与展望
# ═══════════════════════════════════════════════
slide8 = add_blank_slide()
add_section_title(slide8, "5. 总结与展望")

# 左侧：优点
add_textbox(slide8, 0.8, 1.5, 5.5, 0.4, "▎系统优点", font_size=20, color=C_GREEN, bold=True)
pros = [
    "● 完整的数据采集→存储→推理→可视化全链路闭环",
    "● LST-GCN嵌入融合架构，四时域MAE < 4 km/h",
    "● 多时域直接监督策略有效抑制长时域误差堆积",
    "● 管理员/执行者双层角色权限，前后端双重鉴权",
    "● 基于高德API的速度代理策略，低部署成本覆盖11路口",
    "● 支持模拟/真实双模式切换，兼顾演示与实战",
]
add_multiline(slide8, 0.8, 2.1, 5.5, 4.3, pros, font_size=13, color=C_DARK, line_spacing=1.8)

# 右侧：缺点与改进
add_textbox(slide8, 7.0, 1.5, 5.5, 0.4, "▎现有不足与改进方向", font_size=20, color=C_RED, bold=True)
cons = [
    "● 训练数据规模有限（9天采集周期）",
    "    → 延长采集周期，积累更多样本",
    "",
    "● 模型输入特征维度偏低（仅速度）",
    "    → 引入天气、节假日、道路属性等多源特征",
    "",
    "● 11个节点规模较小",
    "    → 扩展至更大范围路网，验证模型泛化性",
    "",
    "● 云端CPU推理耗时高于本地GPU",
    "    → 后续考虑GPU服务器或模型量化加速",
]
add_multiline(slide8, 7.0, 2.1, 5.5, 4.3, cons, font_size=13, color=C_DARK, line_spacing=1.3)

# 底部条
add_rect(slide8, 0, 6.6, 13.333, 0.9, C_PRIMARY)
add_textbox(slide8, 1.0, 6.75, 11.333, 0.5,
    "感谢各位老师！请各位老师批评指正。",
    font_size=24, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_page_number(slide8, 8)

# ── 保存 ──────────────────────────────────────
out_path = os.path.join(BASE, "答辩演示PPT.pptx")
prs.save(out_path)
print(f"\nPPT saved: {out_path}")
print(f"  Total slides: {len(prs.slides)}")
